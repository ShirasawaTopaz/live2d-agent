from internal.agent.tool.base import Tool
from typing import Any, Optional
import os
from internal.agent.sandbox import SandboxMiddleware, default_sandbox


class SandboxedFileToolBase:
    """Base class that adds sandbox checking to file operations."""

    def __init__(self, sandbox: Optional[SandboxMiddleware] = None):
        self.sandbox = sandbox or default_sandbox

    def _check_path(self, path: str, is_write: bool) -> tuple[bool, str]:
        """Check if path access is allowed by sandbox.

        Returns:
            (allowed: bool, error_message: str)
        """
        if not self.sandbox.is_enabled():
            return True, ""

        allowed, reason, normalized = self.sandbox.check_file_access(path, is_write)

        if not allowed:
            if reason.startswith("APPROVAL_REQUIRED"):
                if normalized is not None and self.sandbox.needs_file_approval(
                    normalized, is_write
                ):
                    approved = self.sandbox.request_file_approval(
                        normalized, is_write, reason
                    )
                    if approved:
                        return True, ""
                    else:
                        return False, "Operation was rejected by user approval"
            return False, f"Sandbox denied access: {reason}"

        if not is_write and normalized is not None:
            size_allowed, size_reason = self.sandbox.check_file_size(normalized)
            if not size_allowed:
                return False, size_reason

        return True, ""


class OfficeTool(Tool, SandboxedFileToolBase):
    def __init__(self, sandbox: Optional[SandboxMiddleware] = None):
        SandboxedFileToolBase.__init__(self, sandbox)

    @property
    def name(self) -> str:
        return "office"

    @property
    def description(self) -> str:
        return "Office 文件读写工具，支持读取 docx、xlsx、pptx、pdf 文件内容，也可以写入文本内容到这些文件。当需要处理 Office 文档时调用此工具。"

    @property
    def parameters(self) -> dict:
        return {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["read", "write"],
                    "description": "操作类型：read(读取文件), write(写入文件)",
                },
                "file_path": {
                    "type": "string",
                    "description": "Office 文件路径，支持 .docx .xlsx .pptx .pdf",
                },
                "content": {
                    "type": "string",
                    "description": "写入内容（write 操作时需要）。对于 xlsx，使用表格格式，行用换行分隔，单元格用制表符分隔。",
                },
            },
            "required": ["action", "file_path"],
        }

    async def execute(self, **kwargs) -> Any:
        action = kwargs.get("action")
        file_path = kwargs.get("file_path")

        if not file_path:
            return "Error: File path is required"

        is_write = action == "write"
        allowed, error = self._check_path(file_path, is_write)
        if not allowed:
            return f"Error: {error}"

        if not os.path.exists(file_path) and not is_write:
            return f"Error: File not found: {file_path}"

        ext = os.path.splitext(file_path)[1].lower()

        if action == "read":
            return await self._read_file(file_path, ext)
        elif action == "write":
            content = kwargs.get("content", "")
            return await self._write_file(file_path, ext, content)
        else:
            return f"Error: Unknown action: {action}"

    async def _read_file(self, file_path: str, ext: str) -> Any:
        try:
            if ext == ".docx":
                return self._read_docx(file_path)
            elif ext == ".xlsx":
                return self._read_xlsx(file_path)
            elif ext == ".pptx":
                return self._read_pptx(file_path)
            elif ext == ".pdf":
                return self._read_pdf(file_path)
            else:
                return f"Error: Unsupported file format: {ext}. Supported formats: docx, xlsx, pptx, pdf"
        except ImportError as e:
            return f"Error: Missing dependency: {str(e)}. Please install required packages: pip install python-docx openpyxl python-pptx pymupdf"
        except Exception as e:
            return f"Error reading file: {str(e)}"

    async def _write_file(self, file_path: str, ext: str, content: str) -> Any:
        try:
            if ext == ".docx":
                return self._write_docx(file_path, content)
            elif ext == ".xlsx":
                return self._write_xlsx(file_path, content)
            elif ext == ".pptx":
                return self._write_pptx(file_path, content)
            elif ext == ".pdf":
                return "Error: PDF writing is not supported. PDF format is not suitable for simple text writing."
            else:
                return f"Error: Unsupported file format: {ext}. Supported formats: docx, xlsx, pptx"
        except ImportError as e:
            return f"Error: Missing dependency: {str(e)}. Please install required packages: pip install python-docx openpyxl python-pptx pymupdf"
        except Exception as e:
            return f"Error writing file: {str(e)}"

    def _read_docx(self, file_path: str) -> str:
        from docx import Document

        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs]
        return "\n".join(paragraphs)

    def _write_docx(self, file_path: str, content: str) -> str:
        from docx import Document

        doc = Document()
        for paragraph in content.split("\n"):
            doc.add_paragraph(paragraph)
        doc.save(file_path)
        return "DOCX file written successfully"

    def _read_xlsx(self, file_path: str) -> str:
        from openpyxl import load_workbook

        wb = load_workbook(filename=file_path, read_only=True, data_only=True)
        result = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            result.append(f"=== Sheet: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                row_str = "\t".join(
                    str(cell) if cell is not None else "" for cell in row
                )
                result.append(row_str)
        wb.close()
        return "\n".join(result)

    def _write_xlsx(self, file_path: str, content: str) -> str:
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        for row_idx, line in enumerate(content.splitlines(), start=1):
            cells = line.split("\t")
            for col_idx, cell_value in enumerate(cells, start=1):
                ws.cell(row=row_idx, column=col_idx, value=cell_value.strip())
        wb.save(file_path)
        return "XLSX file written successfully"

    def _read_pptx(self, file_path: str) -> str:
        from pptx import Presentation

        prs = Presentation(file_path)
        result = []
        for idx, slide in enumerate(prs.slides, start=1):
            result.append(f"=== Slide {idx} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    result.append(shape.text)
        return "\n".join(result)

    def _write_pptx(self, file_path: str, content: str) -> str:
        from pptx import Presentation

        prs = Presentation()
        slides_content = content.split("=== Slide")
        for idx, slide_content in enumerate(slides_content):
            if not slide_content.strip():
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title_box = slide.shapes.title
            if idx == 0:
                lines = slide_content.strip().splitlines()
                if lines:
                    title_box.text = lines[0]
                    content_joined = "\n".join(lines[1:])
            else:
                title_box.text = f"Slide {idx}"
                content_joined = slide_content.strip()
            body_box = slide.placeholders[1]
            body_box.text = content_joined
        prs.save(file_path)
        return "PPTX file written successfully"

    def _read_pdf(self, file_path: str) -> str:
        import fitz

        doc = fitz.open(file_path)
        result = []
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text()
            result.append(f"=== Page {page_num} ===")
            result.append(text)
        doc.close()
        return "\n".join(result)
